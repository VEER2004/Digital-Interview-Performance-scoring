from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK_BG = RGBColor(30, 27, 75)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(139, 92, 246)
GREEN = RGBColor(16, 185, 129)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG

def add_title_slide(title_text, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(2.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = subtitle_text
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(148, 163, 184)
    p2.font.name = "Arial"
    p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(title_text, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(22)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.font.name = "Arial"
        p.space_after = Pt(12)
    return slide

def add_image_slide(title_text, img_path, caption=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = "Arial"

    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(1.5), Inches(1.3), height=Inches(5.5))

    if caption:
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(7), Inches(11), Inches(0.4))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(148, 163, 184)
        p2.alignment = PP_ALIGN.CENTER
    return slide

# ============ SLIDES ============

# 1. Title
add_title_slide(
    "🎯 Digital Interview Performance\nScoring System",
    "An AI-Driven Proctoring & Evaluation Platform\n\nSubmitted by: Vir K. Gusai (12202110501063)\nBrainyBeam Info-Tech Pvt. Ltd.\nGCET | CVM University"
)

# 2. Problem Statement
add_content_slide("The Problem", [
    "🔴 Human Bias — Interviewers are prone to unconscious favoritism",
    "🔴 Inconsistency — Evaluation standards vary between interviewers",
    "🔴 Scalability — High-volume recruitment is slow and manual",
    "🔴 Proctoring Integrity — Difficult to verify cheating in remote interviews",
    "🔴 No Data — Traditional interviews produce no quantifiable metrics",
])

# 3. Objective
add_content_slide("Project Objective", [
    "✅ Build a Supervised ML model to predict interview performance scores",
    "✅ Integrate Computer Vision for facial presence & emotion detection",
    "✅ Apply NLP for sentiment analysis on candidate speech transcripts",
    "✅ Implement live WebRTC proctoring with cheat detection",
    "✅ Create an AI Video Scorer for asynchronous evaluation",
    "✅ Generate Power BI-style dashboards with actionable suggestions",
])

# 4. Tech Stack
add_content_slide("Technology Stack", [
    "🐍 Python 3.10+ — Core Runtime",
    "📊 Scikit-Learn — ML Models (Random Forest, Gradient Boosting, Ensemble)",
    "👁️ OpenCV — Haar Cascade face and smile detection",
    "💬 TextBlob — NLP Sentiment Analysis",
    "🎙️ SpeechRecognition + PyAudio — Speech-to-text transcription",
    "🎬 MoviePy — Video-to-audio extraction for AI Video Scorer",
    "🌐 Streamlit + WebRTC — Web UI + Real-time proctoring",
    "📈 Plotly — Interactive gauge charts and radar charts",
])

# 5. Dataset
add_content_slide("Dataset Details", [
    "📁 Source: Kaggle — Virtual Interview Performance Dataset",
    "📊 14 Features: Age, Education Score, Coding Test Score, Confidence,",
    "     Eye Contact, Speech Speed, Filler Words, Interviewer Rating, etc.",
    "🎯 Target: Final Interview Performance Score (Continuous)",
    "🔢 Task Type: Supervised Regression",
    "📐 Preprocessing: StandardScaler normalization, 80/20 train-test split",
])

# 6. ML Pipeline
add_content_slide("Machine Learning Pipeline", [
    "Step 1: Load & preprocess virtual_interview_with_target.csv",
    "Step 2: Train 4 base models — Linear, Random Forest, AdaBoost, GradBoost",
    "Step 3: Build VotingRegressor (Ensemble) combining all 4",
    "Step 4: Evaluate R², MAE, RMSE on 20% test set",
    "Step 5: Select best model (Ensemble if R² ≥ best individual)",
    "Step 6: Save model.pkl, scaler.pkl, feature_columns.pkl via Joblib",
])

# 7. EDA — Correlation Heatmap
add_image_slide("EDA: Correlation Heatmap",
    r"c:\Users\Vir\Desktop\BrainyBEAM Final\eda_plots\correlation_heatmap.png",
    "Feature-to-feature and feature-to-target correlation matrix")

# 8. EDA — Feature Importance
add_image_slide("EDA: Feature Importance",
    r"c:\Users\Vir\Desktop\BrainyBEAM Final\eda_plots\feature_importance.png",
    "Random Forest feature importance ranking")

# 9. EDA — Score Distribution
add_image_slide("EDA: Final Score Distribution",
    r"c:\Users\Vir\Desktop\BrainyBEAM Final\eda_plots\distribution_final_score.png",
    "Histogram of target variable distribution")

# 10. App UI Screenshot
add_image_slide("Application UI — Premium Dark Theme",
    r"C:\Users\Vir\.gemini\antigravity\brain\f1aab91c-ae00-4933-9377-917d749e41d0\initial_load_1778423495960.png",
    "Glassmorphic inputs, gradient buttons, Outfit font, radial gradient background")

# 11. Dashboard Screenshot
add_image_slide("Output: Power BI-Style Dashboard",
    r"C:\Users\Vir\.gemini\antigravity\brain\f1aab91c-ae00-4933-9377-917d749e41d0\final_prediction_dashboard_1778423636709.png",
    "Gauge chart, radar chart, and actionable suggestions")

# 12. Modules Overview
add_content_slide("System Modules Overview", [
    "📋 Tabs 1-4: Manual metric input (14 fields in 2/3-column layouts)",
    "🧠 Tab 5: Advanced Analytics — NLP sentiment, OpenCV face, voice tone",
    "🎥 Tab 6: Live Proctoring — WebRTC video + audio + screen share",
    "     → Real-time face tracking with bounding boxes",
    "     → Live speech transcription via Google STT",
    "     → Tab-switch cheat detection via JavaScript injection",
    "🎬 Tab 7: AI Video Scorer — Upload video → auto-score everything",
])

# 13. Live Proctoring Details
add_content_slide("Live Proctoring System", [
    "👤 Face Tracking: OpenCV detects 0, 1, or multiple faces per frame",
    "     → Green box = candidate detected",
    "     → Red warning = no face OR multiple faces",
    "🎙️ Audio Processing: Buffers mic frames → transcribes on demand",
    "🖥️ Screen Share: JavaScript getDisplayMedia() observes candidate screen",
    "⚠️ Cheat Detection: visibilitychange event fires alert on tab switch",
])

# 14. AI Video Scorer
add_content_slide("AI Video Scorer — Automated Pipeline", [
    "Step 1: MoviePy extracts audio → SpeechRecognition transcribes speech",
    "Step 2: OpenCV analyzes every 15th frame for face + smile detection",
    "Step 3: TextBlob computes NLP sentiment on the full transcript",
    "Step 4: Compute Eye Contact, Confidence, Voice Tone, NLP scores",
    "Step 5: Generate Plotly Gauge + Radar dashboard",
    "Step 6: AI produces verdict — Excellent / Good / Needs Improvement",
])

# 15. Scoring Logic
add_content_slide("Scoring & Categorization Logic", [
    "Manual Mode (ML Prediction):",
    "  → Score ≥ 20: 🌟 Excellent — Fast-track recommended",
    "  → Score ≥ 12: 👍 Good — Consider with improvements",
    "  → Score < 12: ⚠️ Needs Improvement — Not recommended",
    "",
    "AI Video Scorer (CV + NLP):",
    "  → Overall ≥ 7/10: 🌟 Excellent Candidate",
    "  → Overall ≥ 5/10: 👍 Good Candidate",
    "  → Overall < 5/10: ⚠️ Needs Significant Improvement",
])

# 16. Future Enhancements
add_content_slide("Future Enhancements", [
    "🧠 Deep Learning: Replace Haar Cascades with CNN (DeepFace/FER+)",
    "🎙️ Streaming ASR: Use OpenAI Whisper for real-time transcription",
    "💾 Database: PostgreSQL for persistent interview history",
    "📄 PDF Reports: Auto-generate downloadable candidate scorecards",
    "🌍 Multi-language: Extend NLP and STT beyond English",
    "🤸 Body Language: MediaPipe Pose for gesture analysis",
])

# 17. Thank You
add_title_slide(
    "Thank You!",
    "Digital Interview Performance Scoring System\nVir K. Gusai | BrainyBeam Info-Tech Pvt. Ltd.\nGCET | CVM University\n\nQuestions?"
)

# Save
output = r"c:\Users\Vir\Desktop\BrainyBEAM Final\Digital_Interview_Presentation.pptx"
prs.save(output)
print("Presentation saved successfully: " + output)
print("Total slides: " + str(len(prs.slides)))
